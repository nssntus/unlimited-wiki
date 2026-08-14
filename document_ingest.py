"""Local, bounded text extraction for Raw document uploads."""

from __future__ import annotations

import csv
import hashlib
import io
import json
import math
import os
import re
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path

from storage import atomic_write


MAX_INPUT_BYTES = 10 * 1024 * 1024
MAX_EXTRACTED_CHARS = 5_000_000
MAX_ARCHIVE_ENTRIES = 5_000
MAX_ARCHIVE_BYTES = 100 * 1024 * 1024
MAX_IMAGE_PIXELS = 25_000_000

TEXT_SUFFIXES = {
    ".md", ".markdown", ".txt", ".csv", ".tsv", ".json", ".yaml", ".yml", ".xml",
}
HTML_SUFFIXES = {".html", ".htm"}
WORD_SUFFIXES = {".docx"}
SHEET_SUFFIXES = {".xlsx", ".xlsm"}
SLIDE_SUFFIXES = {".pptx", ".pptm"}
CONVERT_WORD_SUFFIXES = {".doc", ".docm", ".odt", ".rtf"}
CONVERT_SHEET_SUFFIXES = {".xls", ".ods"}
CONVERT_SLIDE_SUFFIXES = {".ppt", ".odp"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp", ".gif", ".heic"}
SUPPORTED_SUFFIXES = (
    TEXT_SUFFIXES | HTML_SUFFIXES | WORD_SUFFIXES | SHEET_SUFFIXES | SLIDE_SUFFIXES
    | CONVERT_WORD_SUFFIXES | CONVERT_SHEET_SUFFIXES | CONVERT_SLIDE_SUFFIXES
    | IMAGE_SUFFIXES | {".pdf", ".epub"}
)
PARSER_VERSION = "2"


@dataclass(frozen=True)
class ParsedDocument:
    markdown: str
    source_format: str
    extracted_chars: int
    used_ocr: bool = False


class _VisibleTextParser(HTMLParser):
    BLOCKS = {"article", "blockquote", "br", "div", "h1", "h2", "h3", "h4", "h5", "h6", "li", "p", "section", "td", "th", "tr"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self.ignored = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in {"script", "style", "noscript"}:
            self.ignored += 1
        elif not self.ignored and tag in self.BLOCKS:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript"}:
            self.ignored = max(0, self.ignored - 1)
        elif not self.ignored and tag in self.BLOCKS:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self.ignored:
            self.parts.append(data)

    def text(self) -> str:
        lines = [re.sub(r"\s+", " ", line).strip() for line in "".join(self.parts).splitlines()]
        return "\n\n".join(line for line in lines if line)


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-16", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError("文档文字编码无法识别，未加入原料箱")


def _clean_text(text: str) -> str:
    text = text.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n").replace("\f", "\n\n")
    text = re.sub(r"[ \t]+\n", "\n", text).strip()
    if len(text) > MAX_EXTRACTED_CHARS:
        raise ValueError("文档提取后的文字超过 500 万字符，未加入原料箱")
    return text


def _finish(text: str, suffix: str, *, used_ocr: bool = False) -> ParsedDocument:
    cleaned = _clean_text(text)
    if not re.search(r"[\w\u3400-\u9fff]", cleaned, re.UNICODE):
        if used_ocr:
            raise ValueError("图片中未识别到文字，未加入原料箱")
        raise ValueError("文档中没有可读取的文字，未加入原料箱")
    return ParsedDocument(cleaned + "\n", suffix.removeprefix(".").upper(), len(cleaned), used_ocr)


def _validate_zip(data: bytes) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            infos = archive.infolist()
            total = sum(info.file_size for info in infos)
            if len(infos) > MAX_ARCHIVE_ENTRIES or total > MAX_ARCHIVE_BYTES:
                raise ValueError("文档压缩内容异常，未加入原料箱")
            for info in infos:
                if info.compress_size == 0 and info.file_size > 0:
                    raise ValueError("文档压缩内容异常，未加入原料箱")
                if info.compress_size and info.file_size / info.compress_size > 200:
                    raise ValueError("文档压缩比异常，未加入原料箱")
    except zipfile.BadZipFile as exc:
        raise ValueError("文档文件已损坏或格式不正确，未加入原料箱") from exc


def _markdown_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [[cell.replace("|", "\\|").replace("\n", " ").strip() for cell in row] + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    return "\n".join([
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
        *("| " + " | ".join(row) + " |" for row in normalized[1:]),
    ])


def _parse_csv(data: bytes, delimiter: str) -> str:
    rows = list(csv.reader(io.StringIO(_decode_text(data)), delimiter=delimiter))
    return _markdown_table(rows)


def _parse_html(text: str) -> str:
    parser = _VisibleTextParser()
    parser.feed(text)
    return parser.text()


def _parse_docx(data: bytes) -> str:
    _validate_zip(data)
    try:
        from docx import Document
    except ImportError as exc:
        raise ValueError("Word 解析器未安装，未加入原料箱") from exc
    try:
        document = Document(io.BytesIO(data))
        parts: list[str] = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            match = re.search(r"Heading\s+([1-6])", paragraph.style.name or "", re.I)
            parts.append(("#" * int(match.group(1)) + " " if match else "") + text)
        for table in document.tables:
            rendered = _markdown_table([[cell.text for cell in row.cells] for row in table.rows])
            if rendered:
                parts.append(rendered)
        return "\n\n".join(parts)
    except Exception as exc:
        raise ValueError("Word 文档无法解析，未加入原料箱") from exc


def _parse_xlsx(data: bytes) -> str:
    _validate_zip(data)
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise ValueError("Excel 解析器未安装，未加入原料箱") from exc
    try:
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet in workbook.worksheets:
            rows: list[list[str]] = []
            for index, row in enumerate(sheet.iter_rows(values_only=True)):
                if index >= 20_000:
                    raise ValueError("Excel 单个工作表超过 20000 行，未加入原料箱")
                values = ["" if value is None else str(value) for value in row]
                if any(value.strip() for value in values):
                    while values and not values[-1].strip():
                        values.pop()
                    rows.append(values)
            parts.append(f"## {sheet.title}\n\n{_markdown_table(rows)}")
        workbook.close()
        return "\n\n".join(parts)
    except ValueError as exc:
        if str(exc).startswith("Excel 单个工作表超过"):
            raise
        raise ValueError("Excel 文档无法解析，未加入原料箱") from exc
    except Exception as exc:
        raise ValueError("Excel 文档无法解析，未加入原料箱") from exc


def _parse_pptx(data: bytes) -> str:
    _validate_zip(data)
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError("PowerPoint 解析器未安装，未加入原料箱") from exc
    try:
        presentation = Presentation(io.BytesIO(data))
        parts: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                parts.append(f"## 幻灯片 {index}\n\n" + "\n\n".join(texts))
        return "\n\n".join(parts)
    except Exception as exc:
        raise ValueError("PowerPoint 文档无法解析，未加入原料箱") from exc


def _parse_pdf(data: bytes) -> str:
    try:
        import fitz
    except ImportError as exc:
        raise ValueError("PDF 解析器未安装，未加入原料箱") from exc
    try:
        document = fitz.open(stream=data, filetype="pdf")
        if document.needs_pass:
            raise ValueError("PDF 受密码保护，未加入原料箱")
        if document.page_count > 1_000:
            raise ValueError("PDF 超过 1000 页，未加入原料箱")
        parts = []
        for index, page in enumerate(document, start=1):
            text = page.get_text("text").strip()
            if text:
                parts.append(f"## 第 {index} 页\n\n{text}")
        document.close()
        return "\n\n".join(parts)
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError("PDF 无法解析，未加入原料箱") from exc


def _parse_epub(data: bytes) -> str:
    _validate_zip(data)
    try:
        parts: list[str] = []
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            for name in sorted(archive.namelist()):
                if Path(name).suffix.lower() not in {".html", ".htm", ".xhtml"}:
                    continue
                text = _parse_html(_decode_text(archive.read(name)))
                if text:
                    parts.append(text)
        return "\n\n".join(parts)
    except Exception as exc:
        raise ValueError("EPUB 无法解析，未加入原料箱") from exc


def _ocr_image(filename: str, data: bytes) -> str:
    tesseract = shutil.which("tesseract")
    if not tesseract:
        raise ValueError("图片 OCR 组件未安装，未加入原料箱")
    try:
        from PIL import Image, ImageOps
    except ImportError as exc:
        raise ValueError("图片解析器未安装，未加入原料箱") from exc
    suffix = Path(filename).suffix.lower()
    with tempfile.TemporaryDirectory(prefix="wiki-ocr-") as directory:
        source = Path(directory) / ("source" + suffix)
        source.write_bytes(data)
        image_path = source
        if suffix == ".heic":
            sips = shutil.which("sips")
            if not sips:
                raise ValueError("HEIC 图片解析器不可用，未加入原料箱")
            image_path = Path(directory) / "source.png"
            result = subprocess.run([sips, "-s", "format", "png", str(source), "--out", str(image_path)], capture_output=True, timeout=30)
            if result.returncode != 0:
                raise ValueError("HEIC 图片无法解析，未加入原料箱")
        try:
            with Image.open(image_path) as image:
                image.seek(0)
                image = ImageOps.exif_transpose(image).convert("RGB")
                pixels = image.width * image.height
                if pixels > MAX_IMAGE_PIXELS:
                    ratio = math.sqrt(MAX_IMAGE_PIXELS / pixels)
                    image.thumbnail((max(1, int(image.width * ratio)), max(1, int(image.height * ratio))))
                normalized = Path(directory) / "ocr.png"
                image.save(normalized, format="PNG")
        except Exception as exc:
            raise ValueError("图片文件已损坏或格式不正确，未加入原料箱") from exc
        language = os.environ.get("WIKI_OCR_LANG", "chi_sim+eng")
        try:
            result = subprocess.run(
                [tesseract, str(normalized), "stdout", "-l", language, "--psm", "3"],
                capture_output=True, timeout=60,
            )
        except subprocess.TimeoutExpired as exc:
            raise ValueError("图片 OCR 超时，未加入原料箱") from exc
        if result.returncode != 0:
            raise ValueError("图片 OCR 失败，未加入原料箱")
        return _decode_text(result.stdout)


def _convert_legacy(filename: str, data: bytes, target: str) -> ParsedDocument:
    soffice = shutil.which("soffice")
    if not soffice:
        raise ValueError("旧版 Office/OpenDocument 解析需要 LibreOffice，未加入原料箱")
    with tempfile.TemporaryDirectory(prefix="wiki-document-") as directory:
        root = Path(directory)
        source = root / Path(filename).name
        source.write_bytes(data)
        profile = root / "profile"
        try:
            result = subprocess.run(
                [soffice, "--headless", f"-env:UserInstallation={profile.as_uri()}", "--convert-to", target, "--outdir", str(root), str(source)],
                capture_output=True, timeout=45,
            )
        except subprocess.TimeoutExpired as exc:
            raise ValueError("文档格式转换超时，未加入原料箱") from exc
        converted = next((path for path in root.glob(f"*.{target}") if path != source), None)
        if result.returncode != 0 or converted is None:
            raise ValueError("文档格式转换失败，未加入原料箱")
        parsed = parse_document(converted.name, converted.read_bytes())
        return ParsedDocument(
            parsed.markdown, Path(filename).suffix.removeprefix(".").upper(),
            parsed.extracted_chars, parsed.used_ocr,
        )


def parse_document(filename: str, data: bytes) -> ParsedDocument:
    if not data:
        raise ValueError("Raw 内容不能为空")
    if len(data) > MAX_INPUT_BYTES:
        raise ValueError("Raw 文件超过 10 MiB")
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_SUFFIXES:
        raise ValueError("不支持该文件格式，未加入原料箱")
    if suffix in TEXT_SUFFIXES:
        if suffix == ".csv":
            return _finish(_parse_csv(data, ","), suffix)
        if suffix == ".tsv":
            return _finish(_parse_csv(data, "\t"), suffix)
        return _finish(_decode_text(data), suffix)
    if suffix in HTML_SUFFIXES:
        return _finish(_parse_html(_decode_text(data)), suffix)
    if suffix in WORD_SUFFIXES:
        return _finish(_parse_docx(data), suffix)
    if suffix in SHEET_SUFFIXES:
        return _finish(_parse_xlsx(data), suffix)
    if suffix in SLIDE_SUFFIXES:
        return _finish(_parse_pptx(data), suffix)
    if suffix == ".pdf":
        return _finish(_parse_pdf(data), suffix)
    if suffix == ".epub":
        return _finish(_parse_epub(data), suffix)
    if suffix in IMAGE_SUFFIXES:
        title = Path(filename).stem.strip() or "图片原材料"
        ocr_text = _clean_text(_ocr_image(filename, data))
        meaningful = len(re.findall(r"[\u3400-\u9fff]", ocr_text)) >= 2 or len(re.findall(r"[A-Za-z0-9]", ocr_text)) >= 2
        if not meaningful:
            raise ValueError("图片中未识别到文字，未加入原料箱")
        return _finish(f"# {title}\n\n{ocr_text}", suffix, used_ocr=True)
    if suffix in CONVERT_WORD_SUFFIXES:
        return _convert_legacy(filename, data, "docx")
    if suffix in CONVERT_SHEET_SUFFIXES:
        return _convert_legacy(filename, data, "xlsx")
    if suffix in CONVERT_SLIDE_SUFFIXES:
        return _convert_legacy(filename, data, "pptx")
    raise ValueError("不支持该文件格式，未加入原料箱")


def parse_document_cached(project_root: Path, filename: str, data: bytes) -> ParsedDocument:
    identity = hashlib.sha256(PARSER_VERSION.encode() + b"\0" + filename.encode("utf-8") + b"\0" + data).hexdigest()
    cache_root = project_root / ".wiki-state" / "extracted"
    cache_path = cache_root / f"{identity}.json"
    try:
        payload = json.loads(cache_path.read_text(encoding="utf-8"))
        if payload.get("parser_version") == PARSER_VERSION:
            return ParsedDocument(
                markdown=payload["markdown"], source_format=payload["source_format"],
                extracted_chars=int(payload["extracted_chars"]), used_ocr=bool(payload["used_ocr"]),
            )
    except (FileNotFoundError, ValueError, TypeError, KeyError, json.JSONDecodeError):
        pass
    parsed = parse_document(filename, data)
    cache_root.mkdir(parents=True, exist_ok=True)
    os.chmod(cache_root, 0o700)
    atomic_write(cache_path, json.dumps({
        "parser_version": PARSER_VERSION, "markdown": parsed.markdown,
        "source_format": parsed.source_format, "extracted_chars": parsed.extracted_chars,
        "used_ocr": parsed.used_ocr,
    }, ensure_ascii=False, separators=(",", ":")).encode("utf-8"))
    os.chmod(cache_path, 0o600)
    return parsed
