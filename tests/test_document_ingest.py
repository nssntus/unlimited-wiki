from __future__ import annotations

import io
import zipfile

import fitz
import pytest
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation

import document_ingest
from document_ingest import parse_document
from wiki_service import WikiService


@pytest.fixture
def service(kb_root):
    instance = WikiService(kb_root, start_worker=False)
    try:
        yield instance
    finally:
        instance.close()


def _docx_bytes() -> bytes:
    document = Document()
    document.add_heading("Word 标题", level=1)
    document.add_paragraph("Word 正文内容")
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "数据"
    sheet.append(["名称", "数量"])
    sheet.append(["原材料", 3])
    output = io.BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "演示标题"
    slide.placeholders[1].text = "演示正文内容"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _pdf_bytes() -> bytes:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "PDF readable content")
    data = document.tobytes()
    document.close()
    return data


@pytest.mark.parametrize(
    ("filename", "factory", "expected"),
    [
        ("sample.docx", _docx_bytes, "Word 正文内容"),
        ("sample.xlsx", _xlsx_bytes, "原材料"),
        ("sample.pptx", _pptx_bytes, "演示正文内容"),
        ("sample.pdf", _pdf_bytes, "PDF readable content"),
        ("sample.csv", lambda: "名称,数量\n原材料,3\n".encode(), "| 原材料 | 3 |"),
        ("sample.html", lambda: b"<h1>Title</h1><script>bad()</script><p>Visible text</p>", "Visible text"),
    ],
)
def test_mainstream_documents_extract_text(filename, factory, expected):
    parsed = parse_document(filename, factory())
    assert expected in parsed.markdown
    assert parsed.extracted_chars > 0


def test_binary_raw_preserves_original_and_exposes_text(service):
    data = _docx_bytes()
    uploaded = service.upload_raw("source.docx", data)
    path = uploaded["raw"]["path"]
    assert (service.root / path).read_bytes() == data
    raw = service.read_raw(path)
    assert "Word 正文内容" in raw["markdown"]
    assert raw["source_format"] == "DOCX"
    assert raw["revision"] == uploaded["raw"]["byte_hash"]


def test_image_without_text_is_rejected_before_raw_write(service, monkeypatch):
    image = Image.new("RGB", (240, 120), "white")
    output = io.BytesIO()
    image.save(output, format="PNG")
    monkeypatch.setattr(document_ingest, "_ocr_image", lambda _filename, _data: "")
    before = sorted(path.name for path in (service.root / "raw" / "local").iterdir())
    with pytest.raises(ValueError, match="图片中未识别到文字"):
        service.upload_raw("blank.png", output.getvalue())
    assert sorted(path.name for path in (service.root / "raw" / "local").iterdir()) == before


def test_image_with_ocr_text_is_ingestable(service, monkeypatch):
    image = Image.new("RGB", (240, 120), "white")
    output = io.BytesIO()
    image.save(output, format="PNG")
    monkeypatch.setattr(document_ingest, "_ocr_image", lambda _filename, _data: "图片中的可读文字")
    uploaded = service.upload_raw("screenshot.png", output.getvalue())
    assert uploaded["raw"]["used_ocr"] is True
    assert "图片中的可读文字" in service.read_raw(uploaded["raw"]["path"])["markdown"]


def test_image_with_only_one_ocr_noise_character_is_rejected(service, monkeypatch):
    image = Image.new("RGB", (240, 120), "white")
    output = io.BytesIO()
    image.save(output, format="PNG")
    monkeypatch.setattr(document_ingest, "_ocr_image", lambda _filename, _data: "x")
    with pytest.raises(ValueError, match="图片中未识别到文字"):
        service.upload_raw("noise.png", output.getvalue())
    assert not (service.root / "raw" / "local" / "noise.png").exists()


def test_image_ocr_result_is_cached_per_workspace(service, monkeypatch):
    image = Image.new("RGB", (240, 120), "white")
    output = io.BytesIO()
    image.save(output, format="PNG")
    calls = []
    monkeypatch.setattr(document_ingest, "_ocr_image", lambda _filename, _data: calls.append(True) or "缓存文字")
    service.upload_raw("cached.png", output.getvalue())
    after_upload = len(calls)
    service.raw_inbox()
    service.read_raw("raw/local/cached.png")
    assert after_upload == 1
    assert len(calls) == after_upload
    cache_files = list((service.root / ".wiki-state" / "extracted").glob("*.json"))
    assert cache_files and all(path.stat().st_mode & 0o077 == 0 for path in cache_files)


def test_archive_bomb_is_rejected():
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", "A" * 1_000_000)
    with pytest.raises(ValueError, match="压缩比异常"):
        parse_document("bomb.docx", output.getvalue())


@pytest.mark.parametrize(
    ("filename", "content_type"),
    [("broken.docx", "Word"), ("broken.xlsx", "Excel"), ("broken.pptx", "PowerPoint")],
)
def test_malformed_office_files_return_controlled_error(filename, content_type):
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w") as archive:
        archive.writestr("[Content_Types].xml", "<not-valid")
    with pytest.raises(ValueError, match=rf"{content_type} 文档无法解析"):
        parse_document(filename, output.getvalue())
